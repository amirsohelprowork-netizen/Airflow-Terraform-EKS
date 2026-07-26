"""Generate an editable PowerPoint overview for this Airflow EC2 project."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "Airflow_EC2_CICD_Setup_and_DAG_Deployment.pptx"

NAVY = RGBColor(12, 30, 55)
BLUE = RGBColor(31, 119, 180)
TEAL = RGBColor(0, 166, 153)
ORANGE = RGBColor(241, 142, 43)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(35, 48, 65)
MUTED = RGBColor(95, 110, 125)
PALE = RGBColor(240, 246, 250)
GREEN = RGBColor(37, 152, 92)
RED = RGBColor(204, 73, 73)


def set_text(shape, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT,
             font="Aptos", margin=0.08):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_box(slide, x, y, w, h, text="", fill=WHITE, line=None, radius=True,
            size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line if line else fill
    if text:
        set_text(box, text, size=size, color=color, bold=bold, align=align)
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def add_title(slide, title, subtitle=None, number=None):
    add_box(slide, 0, 0, 13.333, 0.25, fill=TEAL, radius=False)
    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.46), Inches(11.8), Inches(0.48))
    set_text(title_box, title, size=27, color=NAVY, bold=True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.64), Inches(0.96), Inches(11.8), Inches(0.32))
        set_text(sub, subtitle, size=11, color=MUTED)
    if number is not None:
        num = slide.shapes.add_textbox(Inches(12.15), Inches(7.07), Inches(0.55), Inches(0.25))
        set_text(num, str(number), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide):
    foot = slide.shapes.add_textbox(Inches(0.64), Inches(7.08), Inches(8.5), Inches(0.22))
    set_text(foot, "Airflow on AWS EC2 | Terraform + Docker Compose + GitHub Actions", size=9, color=MUTED)


def bullets(slide, items, x=0.8, y=1.55, w=5.7, h=4.9, size=19):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(14)
        p.bullet = True
    return box


def code_box(slide, code, x, y, w, h, label=None):
    if label:
        add_box(slide, x, y - 0.3, 2.0, 0.28, label, fill=ORANGE, size=10, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)
    box = add_box(slide, x, y, w, h, fill=RGBColor(24, 39, 55), line=RGBColor(24, 39, 55), radius=False)
    set_text(box, code, size=12, color=RGBColor(225, 235, 242), font="Cascadia Mono", margin=0.17)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    return box


def flow(slide, labels, x=0.75, y=3.1, width=2.25):
    for i, (label, color) in enumerate(labels):
        left = x + i * (width + 0.28)
        add_box(slide, left, y, width, 0.88, label, fill=color, size=15, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + width + 0.05), Inches(y + 0.27), Inches(0.18), Inches(0.33))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED; arrow.line.color.rgb = MUTED


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, 13.333, 7.5, fill=NAVY, radius=False)
    add_box(slide, 0, 0, 13.333, 0.22, fill=TEAL, radius=False)
    add_box(slide, 0.76, 1.03, 0.86, 0.86, "A", fill=TEAL, size=35, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    t = slide.shapes.add_textbox(Inches(0.78), Inches(2.08), Inches(11.4), Inches(1.05))
    set_text(t, "Airflow on AWS EC2", size=40, color=WHITE, bold=True)
    s = slide.shapes.add_textbox(Inches(0.82), Inches(3.15), Inches(11), Inches(0.6))
    set_text(s, "Step-by-step setup with Terraform, Docker Compose, and GitHub Actions", size=21, color=RGBColor(196, 220, 235))
    b = slide.shapes.add_textbox(Inches(0.82), Inches(5.65), Inches(8), Inches(0.35))
    set_text(b, "From an empty AWS account to a deployed DAG", size=16, color=RGBColor(147, 205, 197), bold=True)

    # 2. Goal
    slide = prs.slides.add_slide(blank); add_title(slide, "1. What this project delivers", "A lightweight, demo-friendly Airflow platform on AWS", 2); add_footer(slide)
    bullets(slide, ["Infrastructure as code: Terraform builds the AWS networking, security group, EC2 host, Elastic IP, and SSH key.",
                    "Runtime: Docker Compose runs Airflow 2.9.2 with LocalExecutor and PostgreSQL 15.",
                    "Automation: GitHub Actions validates DAG Python files and deploys them to the EC2 host over SSH.",
                    "Outcome: open the Airflow UI, enable a DAG, trigger it, and view task logs."], x=0.75, y=1.55, w=6.7, size=18)
    add_box(slide, 8.1, 1.75, 4.1, 3.7, fill=PALE, line=RGBColor(205, 220, 230))
    set_text(slide.shapes[-1], "Confirmed environment\n\n• Airflow UI reachable\n• example_hello_world parsed\n• LocalExecutor + PostgreSQL\n• Daily sample schedule\n\nUse this as a test/demo platform; harden access before production.", size=18, color=INK, bold=False, margin=0.25)

    # 3 architecture
    slide = prs.slides.add_slide(blank); add_title(slide, "2. Architecture at a glance", "The CI/CD path and the runtime path meet at the Airflow host", 3); add_footer(slide)
    flow(slide, [("Developer\nGit push", BLUE), ("GitHub Actions\nvalidate + copy", TEAL), ("AWS EC2\nDocker Compose", ORANGE), ("Airflow UI\nDAG runs", NAVY)], x=0.55, y=2.75, width=2.8)
    add_box(slide, 1.05, 1.48, 11.2, 0.62, "Terraform provisions the AWS foundation once; GitHub Actions handles repeat DAG deployments.", fill=PALE, line=PALE, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 3.55, 4.45, 6.2, 1.1, "Inside EC2: 4 GB swap | Docker | PostgreSQL metadata DB | Airflow webserver | Airflow scheduler", fill=RGBColor(235, 248, 246), line=TEAL, size=17, color=INK, align=PP_ALIGN.CENTER)

    # 4 prerequisites
    slide = prs.slides.add_slide(blank); add_title(slide, "3. Prerequisites", "Prepare access and tools before provisioning", 4); add_footer(slide)
    bullets(slide, ["AWS account with permission to create EC2, VPC, Elastic IP, security groups, and key pairs.",
                    "AWS CLI configured locally: aws configure (use a least-privilege IAM identity).",
                    "Terraform 1.5+ installed locally for initial provisioning and lifecycle operations.",
                    "GitHub repository with Actions enabled; you need permission to add repository secrets.",
                    "SSH client and a browser. Restrict allowed_cidr_blocks to your public IP where possible."], x=0.75, y=1.5, w=6.7, size=17)
    code_box(slide, "aws configure\nterraform -version\ngit --version\nssh -V", 8.0, 2.05, 3.95, 1.7, "CHECK LOCAL TOOLS")
    add_box(slide, 8.0, 4.25, 3.95, 1.0, "Security note\nNever commit keys, passwords, terraform.tfvars, or terraform.tfstate.", fill=RGBColor(255, 242, 240), line=RED, size=15, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 5 repo
    slide = prs.slides.add_slide(blank); add_title(slide, "4. Understand the repository", "These are the working areas you will use", 5); add_footer(slide)
    code_box(slide, "terraform/\n  main.tf, vpc.tf, ec2.tf\n  security_groups.tf\n  user_data.sh\n  terraform.tfvars\n\ndags/\n  example_dag.py\n\n.github/workflows/\n  deploy-dags.yml\n  terraform.yml\n\nscripts/validate_dags.py", 0.9, 1.55, 4.3, 4.55, "PROJECT LAYOUT")
    bullets(slide, ["terraform/: AWS resource definitions and EC2 bootstrap script.", "dags/: version-controlled Airflow workflows.", "deploy-dags.yml: validates and copies DAG files to /opt/airflow/dags/.", "validate_dags.py: syntax check used before CI/CD deployment."], x=6.0, y=1.72, w=6.1, size=19)

    # 6 configure
    slide = prs.slides.add_slide(blank); add_title(slide, "5. Configure Terraform for your account", "Set safe values before running apply", 6); add_footer(slide)
    bullets(slide, ["Copy terraform.tfvars.example to terraform.tfvars.", "Choose us-east-1 (or your target region) and a suitable instance type. The recorded deployment uses t3.small.", "Use a strong Airflow admin password; do not use the sample password in the repository.", "Set allowed_cidr_blocks to your own public-IP CIDR instead of 0.0.0.0/0 for real use."], x=0.75, y=1.52, w=6.1, size=18)
    code_box(slide, "# terraform/terraform.tfvars\naws_region               = \"us-east-1\"\ninstance_type            = \"t3.small\"\nairflow_image_tag        = \"2.9.2\"\nairflow_admin_username   = \"admin\"\nairflow_admin_password   = \"<strong-secret>\"\nallowed_cidr_blocks      = [\"YOUR.PUBLIC.IP/32\"]", 7.15, 1.55, 5.35, 3.65, "EXAMPLE CONFIGURATION")

    # 7 provision
    slide = prs.slides.add_slide(blank); add_title(slide, "6. Provision the AWS environment", "Terraform creates the host and starts its bootstrap sequence", 7); add_footer(slide)
    code_box(slide, "cd terraform\nterraform init\nterraform fmt -check\nterraform validate\nterraform plan\nterraform apply", 0.85, 1.65, 4.05, 2.4, "RUN LOCALLY")
    bullets(slide, ["Review the plan carefully before approving apply.", "Terraform provisions a VPC, public subnet, internet gateway, security group, EC2 instance, and Elastic IP.", "The private key is saved locally as terraform/airflow-key.pem; protect it and do not add it to Git.", "The EC2 user-data script configures swap, installs Docker/Compose, writes the Compose stack, initializes Airflow, and starts services."], x=5.5, y=1.55, w=6.7, size=17)
    add_box(slide, 0.9, 4.75, 11.35, 0.72, "Wait a few minutes after apply: the bootstrap script downloads images and Airflow services need time to become healthy.", fill=RGBColor(255, 247, 232), line=ORANGE, size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # 8 verify
    slide = prs.slides.add_slide(blank); add_title(slide, "7. Verify the host and open Airflow", "Use Terraform outputs, then confirm the Docker services", 8); add_footer(slide)
    code_box(slide, "cd terraform\nterraform output airflow_url\nterraform output -raw ssh_command\n\n# From your local terminal\nssh -i airflow-key.pem ec2-user@<PUBLIC_IP>\nsudo docker-compose -f /opt/airflow/docker-compose.yml ps", 0.8, 1.55, 5.7, 3.5, "VERIFY")
    bullets(slide, ["Open http://<PUBLIC_IP>:8080 in a browser.", "Sign in with the Airflow admin credentials you configured.", "Expect PostgreSQL, airflow-webserver, and airflow-scheduler to be healthy.", "The sample example_hello_world DAG should appear; it is initially paused by configuration."], x=7.1, y=1.65, w=5.3, size=17)

    # 9 GitHub
    slide = prs.slides.add_slide(blank); add_title(slide, "8. Configure GitHub Actions", "Add secrets once so pushes can deploy DAGs", 9); add_footer(slide)
    add_box(slide, 0.82, 1.5, 5.6, 4.4, fill=PALE, line=RGBColor(205, 220, 230))
    set_text(slide.shapes[-1], "Repository Settings → Secrets and variables → Actions\n\nRequired for DAG deployments\n\nAIRFLOW_HOST\n  EC2 Elastic IP or DNS name\n\nSSH_PRIVATE_KEY\n  Full contents of terraform/airflow-key.pem\n\nKeep both as repository secrets — never add them to YAML or source files.", size=18, color=INK, margin=0.27)
    bullets(slide, ["The workflow triggers on pushes to main that change dags/ or requirements/.", "It installs Airflow 2.9.2, validates DAG syntax, then uses SSH/SCP to copy files to the host.", "Pull requests validate only; deployment runs only after a push to main.", "For Terraform workflow use, configure AWS_ACCESS_KEY_ID, AWS_SECRET_ACCESS_KEY, and AWS_REGION too."], x=7.05, y=1.5, w=5.3, size=17)

    # 10 dag
    slide = prs.slides.add_slide(blank); add_title(slide, "9. Create a DAG", "Put one Python workflow file in the dags/ folder", 10); add_footer(slide)
    code_box(slide, "from datetime import datetime\nfrom airflow import DAG\nfrom airflow.operators.python import PythonOperator\n\ndef greet():\n    print(\"Hello from my deployed DAG\")\n\nwith DAG(\n    dag_id=\"my_first_pipeline\",\n    start_date=datetime(2026, 1, 1),\n    schedule=None, catchup=False,\n) as dag:\n    PythonOperator(task_id=\"greet\", python_callable=greet)", 0.7, 1.38, 6.15, 4.95, "dags/my_first_pipeline.py")
    bullets(slide, ["Give every DAG a unique dag_id.", "Set catchup=False for a simple new pipeline unless backfills are intentional.", "Keep dependencies compatible with the Airflow image; add any needed packages to requirements/requirements.txt.", "Run a local syntax check before committing: python scripts/validate_dags.py."], x=7.25, y=1.65, w=5.25, size=17)

    # 11 deploy
    slide = prs.slides.add_slide(blank); add_title(slide, "10. Deploy a DAG through CI/CD", "This is the normal repeatable deployment path", 11); add_footer(slide)
    flow(slide, [("Edit DAG\nin dags/", BLUE), ("Commit + push\nto main", NAVY), ("Actions\nvalidates", TEAL), ("SSH/SCP\ncopies DAG", ORANGE), ("Scheduler\nloads it", GREEN)], x=0.35, y=2.05, width=2.35)
    code_box(slide, "git add dags/my_first_pipeline.py\ngit commit -m \"feat: add first pipeline\"\ngit push origin main", 3.65, 4.15, 5.95, 1.35, "COMMANDS")
    add_box(slide, 1.6, 5.95, 10.2, 0.58, "Watch the “Deploy DAGs to Airflow” workflow in the GitHub Actions tab. The scheduler can take up to ~60 seconds to discover the new file.", fill=PALE, line=PALE, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # 12 operate
    slide = prs.slides.add_slide(blank); add_title(slide, "11. Run and monitor the deployed DAG", "The UI is where you enable, trigger, inspect, and troubleshoot workflows", 12); add_footer(slide)
    bullets(slide, ["Open the Airflow UI and locate my_first_pipeline.", "Toggle the DAG from paused to active, then click the play button to trigger it manually.", "Open Grid or Graph view to inspect task state; select a task and view its logs.", "For host-side checks: SSH in and run docker-compose -f /opt/airflow/docker-compose.yml ps.", "If a DAG is missing, first inspect the GitHub Actions run, then validate the file and verify its location under /opt/airflow/dags/."], x=0.75, y=1.5, w=7.1, size=18)
    add_box(slide, 8.5, 1.75, 3.45, 3.5, "Healthy signals\n\n✓ UI opens\n✓ DAG appears\n✓ Scheduler is healthy\n✓ Task logs are written\n✓ GitHub Actions run is green", fill=RGBColor(235, 248, 246), line=TEAL, size=19, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # 13 operations
    slide = prs.slides.add_slide(blank); add_title(slide, "12. Operational notes and next improvements", "Keep the demo dependable and safer", 13); add_footer(slide)
    bullets(slide, ["Documentation cleanup: update README.md — it still describes the retired MWAA/S3 design; context.md is closer to the EC2 implementation.",
                    "Security: replace sample credentials, restrict SSH/UI CIDRs, consider HTTPS + a reverse proxy, and rotate exposed credentials.",
                    "Reliability: back up Terraform state remotely, monitor disk/memory, and use a larger instance or managed service for production workloads.",
                    "Dependencies: use a controlled requirements file and test DAG imports in CI before deployment.",
                    "Cost control: when the demo is no longer needed, run terraform destroy from terraform/."], x=0.75, y=1.45, w=11.7, size=17)
    add_box(slide, 1.2, 5.72, 10.9, 0.64, "Final result: infrastructure is reproducible, DAG changes are version-controlled, and deployment is automated by GitHub Actions.", fill=NAVY, line=NAVY, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
